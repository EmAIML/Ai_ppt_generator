# utils/ppt_utils.py
from pptx import Presentation
from pptx.util import Inches
import os, uuid
import logging

log = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

OUTPUT_DIR = "output"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def create_ppt(slides):
    """
    slides: list of dicts like {'title': '...', 'bullets': [...], 'image_path': '/tmp/x.png' or None}
    Returns path to saved pptx.
    """
    prs = Presentation()

    for s in slides:
        # Ensure layout exists; layout 1 typically has title + content. If not, fallback to blank.
        try:
            slide_layout = prs.slide_layouts[1]
        except Exception:
            slide_layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(slide_layout)

        # Some templates may not have a title placeholder in layout 1; guard carefully
        try:
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = s.get('title', '')
        except Exception:
            log.debug("No title placeholder on this layout")

        # Body placeholder index may vary; attempt to find the first text frame placeholder
        body_text = "\n".join(s.get('bullets', []))
        inserted = False
        for shape in slide.placeholders:
            try:
                if shape.has_text_frame:
                    shape.text = body_text
                    inserted = True
                    break
            except Exception:
                continue
        if not inserted:
            # fallback: add a textbox
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(6)
            height = Inches(3)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = body_text

        # add image if available
        img_path = s.get('image_path')
        if img_path:
            try:
                # place on right side; adjust values for your template
                slide.shapes.add_picture(img_path, Inches(5), Inches(1.5), width=Inches(4))
            except Exception as e:
                log.exception("Failed to add picture %s: %s", img_path, e)

    file_path = os.path.join(OUTPUT_DIR, f"presentation_{uuid.uuid4().hex}.pptx")
    prs.save(file_path)
    return file_path
